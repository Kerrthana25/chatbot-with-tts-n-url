import os
import PyPDF2
from docx import Document
import pandas as pd
from pptx import Presentation
import pytesseract
from pdf2image import convert_from_path
from PIL import Image
import cv2
import requests
from bs4 import BeautifulSoup
from urllib.parse import urlparse, urljoin

import re


class DocumentProcessor:
    @staticmethod
    def extract_text_from_pdf(pdf_path):
        text = ""
        with open(pdf_path, 'rb') as file:
            reader = PyPDF2.PdfReader(file)
            for page in reader.pages:
                text += page.extract_text()
        return text

    @staticmethod
    def extract_text_from_docx(docx_path):
        doc = Document(docx_path)
        return '\n'.join([para.text for para in doc.paragraphs])

    @staticmethod
    def extract_text_from_xlsx(xlsx_path):
        try:
            df = pd.read_excel(xlsx_path)
            return df.to_string()
        except:
            return "Could not read Excel file"

    @staticmethod
    def extract_text_from_pptx(pptx_path):
        prs = Presentation(pptx_path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return '\n'.join(text)

    @staticmethod
    def extract_text_from_image(image_path):
        try:
            import cv2
            import pytesseract

            if not os.path.exists(image_path):
                raise FileNotFoundError(f"Image not found: {image_path}")

            img = cv2.imread(image_path)
            if img is None:
                return "Could not open image"

            # Convert to grayscale
            gray = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)

            # Resize image for better accuracy
            gray = cv2.resize(gray, None, fx=1.5, fy=1.5, interpolation=cv2.INTER_LINEAR)

            # Thresholding
            _, thresh = cv2.threshold(gray, 150, 255, cv2.THRESH_BINARY + cv2.THRESH_OTSU)

            # OCR
            config = r'--oem 3 --psm 6'
            text = pytesseract.image_to_string(thresh, config=config).strip()

            return text if text else "No text found in image"

        except Exception as e:
            return f"Error processing image: {str(e)}"

    @staticmethod
    def extract_text_from_video(video_path):
        # Basic video text extraction (from frames)
        try:
            cap = cv2.VideoCapture(video_path)
            text = ""
            while cap.isOpened():
                ret, frame = cap.read()
                if not ret:
                    break
                gray = cv2.cvtColor(frame, cv2.COLOR_BGR2GRAY)
                text += pytesseract.image_to_string(gray)
            cap.release()
            return text
        except:
            return "Could not process video"

    @staticmethod
    def extract_text_from_website(url):
        try:
            response = requests.get(url)
            soup = BeautifulSoup(response.text, 'html.parser')
            for script in soup(["script", "style"]):
                script.decompose()
            return ' '.join(soup.stripped_strings)
        except Exception as e:
            print(f"Error extracting text from website: {e}")
            return ""

    @staticmethod
    def process_file(file_path):
        if file_path.startswith('http://') or file_path.startswith('https://'):
            return DocumentProcessor.extract_text_from_website(file_path)

        extension = os.path.splitext(file_path)[1].lower()[1:]

        processors = {
            'pdf': DocumentProcessor.extract_text_from_pdf,
            'docx': DocumentProcessor.extract_text_from_docx,
            'xlsx': DocumentProcessor.extract_text_from_xlsx,
            'pptx': DocumentProcessor.extract_text_from_pptx,
            'jpg': DocumentProcessor.extract_text_from_image,
            'jpeg': DocumentProcessor.extract_text_from_image,
            'png': DocumentProcessor.extract_text_from_image,
            'mp4': DocumentProcessor.extract_text_from_video,
            'txt': lambda path: open(path, 'r', encoding='utf-8').read(),
            'csv': DocumentProcessor.extract_text_from_xlsx
        }

        if extension in processors:
            try:
                return processors[extension](file_path)
            except Exception as e:
                return f"Error processing {extension} file: {str(e)}"
        else:
            return f"Unsupported file type: {extension}"

    @staticmethod
    def crawl_website(url, visited=None, depth=2):
        """Recursively crawl internal links up to a specified depth."""
        if visited is None:
            visited = set()
        if depth == 0 or url in visited:
            return ""

        print(f"Crawling: {url}")
        visited.add(url)
        try:
            response = requests.get(url, timeout=10, headers={
                'User-Agent': 'Mozilla/5.0'
            })
            soup = BeautifulSoup(response.text, 'html.parser')

            # Remove unwanted tags
            for tag in soup(['script', 'style', 'nav', 'footer', 'header', 'iframe']):
                tag.decompose()

            page_text = soup.get_text(separator='\n', strip=True)

            # Find internal links
            parsed_url = urlparse(url)
            base_url = f"{parsed_url.scheme}://{parsed_url.netloc}"

            links = set()
            for a_tag in soup.find_all('a', href=True):
                link = a_tag['href']
                if link.startswith('/'):
                    link = base_url + link
                elif not link.startswith('http'):
                    continue  # Skip mailto, javascript, etc.

                if base_url in link:
                    links.add(link)

            # Recursively crawl internal links
            for link in links:
                page_text += "\n" + DocumentProcessor.crawl_website(link, visited, depth - 1)

            return page_text

        except Exception as e:
            print(f"Failed to crawl {url}: {e}")
            return ""
